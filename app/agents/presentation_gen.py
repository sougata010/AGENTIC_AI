import os
import io
import random
import requests
import asyncio
from datetime import datetime
from typing import List, Optional, Dict, Any
from pydantic import BaseModel, Field
from langchain_core.prompts import ChatPromptTemplate

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.agents.base import BaseAgent
from app.config import settings


class ColorPalette(BaseModel):
    primary: str = Field(description="Primary color hex")
    secondary: str = Field(description="Panel background color hex")
    accent: str = Field(description="Highlight color hex")
    background: str = Field(description="Slide background color hex")
    text_on_primary: str = Field(description="Text color ON primary elements")
    text_on_secondary: str = Field(description="Text color ON secondary panels")
    text_on_background: str = Field(description="Text color ON main background")

class Fonts(BaseModel):
    heading: str
    body: str

class Theme(BaseModel):
    color_palette: ColorPalette
    fonts: Fonts

class Layout(BaseModel):
    layout_family: str = Field(description="hero_focus, split_visual_text, card_grid, timeline_flow, centered_statement, comparison_columns")
    visual_anchor: str
    structure: List[str]

class Hierarchy(BaseModel):
    level_1: str
    level_2: str
    level_3: str

class Content(BaseModel):
    title: str = Field(description="Short title max 6-8 words")
    subtitle: Optional[str] = Field(description="Short subtitle max 10 words")
    points: List[str] = Field(description="Bullet points, each max 15 words")
    image_keyword: Optional[str] = Field(description="ONE single word for image search like: milkshake, mango, handshake, office")

class AnimationStep(BaseModel):
    element: str
    animation: str
    delay_ms: int

class Emphasis(BaseModel):
    element: str
    style: str

class Motion(BaseModel):
    enter: str
    sequence: List[AnimationStep]
    emphasis: Optional[Emphasis]

class SlideSpec(BaseModel):
    slide_id: int
    slide_type: str
    goal: str
    layout: Layout
    hierarchy: Hierarchy
    content: Content
    motion: Motion

class PresentationSpec(BaseModel):
    theme: Theme
    slides: List[SlideSpec]


DESIGN_SYSTEM_PROMPT = '''
You are a professional presentation designer.

SLIDE TYPES: hero, concept, comparison, process, list_grid, problem, solution, vision
LAYOUTS: hero_focus, split_visual_text, card_grid, timeline_flow, centered_statement, comparison_columns

MAPPING:
hero -> hero_focus
concept -> split_visual_text
process -> timeline_flow
list_grid -> card_grid
comparison -> comparison_columns
vision/problem/solution -> centered_statement

COLOR RULES:
- Choose colors based on TOPIC CONTEXT (food=warm/appetizing, tech=cool/modern, nature=green, etc.)
- Ensure HIGH CONTRAST between text and backgrounds
- text_on_primary/secondary/background must be VISIBLE on their respective backgrounds

CONTENT LENGTH RULES (CRITICAL):
- Titles: MAX 6-8 words
- Subtitles: MAX 10 words
- Bullet points: MAX 15 words each
- Card points: Use "Title: Description" format, keep short
- Timeline points: MAX 8 words each

IMAGE KEYWORD RULES (VERY IMPORTANT):
- Use ONLY ONE SINGLE WORD - the most important concrete noun
- Good examples: "milkshake", "mango", "handshake", "office", "coffee", "laptop"
- NEVER use multiple words or phrases
- NEVER use abstract concepts like "quality", "success", "innovation", "partnership"
- The keyword must be a REAL photographable physical object
'''


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def fetch_image(keyword: str, w: int = 1920, h: int = 1080):
    headers = {'User-Agent': 'Mozilla/5.0'}
    kw = keyword.lower().split()[0].split(',')[0].strip()
    if not kw or len(kw) < 2:
        kw = "abstract"
    
    print(f"    Searching: {kw}")
    
    try:
        url = f"https://source.unsplash.com/{w}x{h}/?{kw}"
        r = requests.get(url, timeout=20, headers=headers, allow_redirects=True)
        if r.status_code == 200 and len(r.content) > 10000:
            return r.content
    except:
        pass
    
    try:
        seed = random.randint(1, 10000)
        url = f"https://loremflickr.com/{w}/{h}/{kw}?lock={seed}"
        r = requests.get(url, timeout=15, headers=headers, allow_redirects=True)
        if r.status_code == 200 and len(r.content) > 5000:
            return r.content
    except:
        pass
    
    try:
        url = f"https://picsum.photos/seed/{random.randint(1,10000)}/{w}/{h}"
        r = requests.get(url, timeout=15, headers=headers, allow_redirects=True)
        if r.status_code == 200:
            return r.content
    except:
        pass
    return None

def add_bg(slide, prs, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def fit_text(text, max_chars, base_size, min_size):
    if len(text) <= max_chars:
        return base_size
    ratio = max_chars / len(text)
    return max(min_size, int(base_size * min(1, ratio * 1.2)))

def truncate(text, max_len):
    if len(text) <= max_len:
        return text
    return text[:max_len-3] + "..."


def render_hero_focus(slide, prs, spec, theme):
    bg = hex_to_rgb(theme.color_palette.background)
    accent = hex_to_rgb(theme.color_palette.accent)
    add_bg(slide, prs, bg)
    
    kw = spec.content.image_keyword or "abstract"
    img = fetch_image(kw, 1920, 1080)
    if img:
        pic = slide.shapes.add_picture(io.BytesIO(img), 0, 0, prs.slide_width, prs.slide_height)
        slide.shapes._spTree.insert(2, pic._element)
        ov = add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RGBColor(0,0,0))
        ov.fill.fore_color.brightness = -0.55
    
    add_rect(slide, Inches(0.8), Inches(3.4), Inches(2), Pt(4), accent)
    
    title = truncate(spec.content.title, 50)
    size = fit_text(title, 30, 54, 32)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.name = theme.fonts.heading
    
    if spec.content.subtitle:
        sub = truncate(spec.content.subtitle, 60)
        sb = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10), Inches(0.8))
        p = sb.text_frame.paragraphs[0]
        p.text = sub
        p.font.size = Pt(22)
        p.font.color.rgb = accent

def render_split_visual_text(slide, prs, spec, theme):
    bg = hex_to_rgb(theme.color_palette.background)
    accent = hex_to_rgb(theme.color_palette.accent)
    secondary = hex_to_rgb(theme.color_palette.secondary)
    txt = hex_to_rgb(theme.color_palette.text_on_secondary)
    add_bg(slide, prs, bg)
    
    kw = spec.content.image_keyword or "office"
    img = fetch_image(kw, 960, 1080)
    if img:
        slide.shapes.add_picture(io.BytesIO(img), Inches(6.66), 0, Inches(6.66), prs.slide_height)
    
    panel = add_rect(slide, 0, 0, Inches(7), prs.slide_height, secondary)
    slide.shapes._spTree.insert(2, panel._element)
    
    add_rect(slide, Inches(0.5), Inches(0.6), Pt(4), Inches(0.8), accent)
    
    title = truncate(spec.content.title, 40)
    size = fit_text(title, 25, 36, 22)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(5.8), Inches(1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = txt
    
    points = spec.content.points[:5]
    bsize = 18 if len(points) <= 4 else 15
    bb = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(5.8), Inches(5))
    tf = bb.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"-> {truncate(pt, 80)}"
        p.font.size = Pt(bsize)
        p.font.color.rgb = txt
        p.space_before = Pt(12)

def render_card_grid(slide, prs, spec, theme):
    bg = hex_to_rgb(theme.color_palette.background)
    accent = hex_to_rgb(theme.color_palette.accent)
    secondary = hex_to_rgb(theme.color_palette.secondary)
    bg_txt = hex_to_rgb(theme.color_palette.text_on_background)
    card_txt = hex_to_rgb(theme.color_palette.text_on_secondary)
    add_bg(slide, prs, bg)
    
    title = truncate(spec.content.title, 40)
    size = fit_text(title, 30, 34, 24)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = bg_txt
    
    add_rect(slide, Inches(0.6), Inches(1.3), Inches(12), Pt(2), accent)
    
    cards = spec.content.points[:4]
    cw = Inches(2.9)
    ch = Inches(4.2)
    gap = Inches(0.25)
    
    for i, pt in enumerate(cards):
        x = Inches(0.6) + i * (cw + gap)
        add_rect(slide, x, Inches(1.6), cw, ch, secondary)
        
        parts = pt.split(':') if ':' in pt else [f"Point {i+1}", pt]
        ct = truncate(parts[0].strip(), 25)
        cb = truncate(parts[1].strip() if len(parts) > 1 else pt, 100)
        
        ttb = slide.shapes.add_textbox(x + Inches(0.15), Inches(1.8), cw - Inches(0.3), Inches(0.7))
        tf = ttb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ct
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent
        
        btb = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.6), cw - Inches(0.3), Inches(3))
        tf = btb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cb
        p.font.size = Pt(12)
        p.font.color.rgb = card_txt

def render_timeline_flow(slide, prs, spec, theme):
    bg = hex_to_rgb(theme.color_palette.background)
    accent = hex_to_rgb(theme.color_palette.accent)
    txt = hex_to_rgb(theme.color_palette.text_on_background)
    add_bg(slide, prs, bg)
    
    title = truncate(spec.content.title, 40)
    size = fit_text(title, 30, 32, 22)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = txt
    
    line_y = Inches(3)
    add_rect(slide, Inches(0.8), line_y, Inches(11.5), Pt(3), accent)
    
    steps = spec.content.points[:5]
    n = len(steps)
    sw = Inches(11.5) / n
    
    for i, pt in enumerate(steps):
        x = Inches(0.8) + i * sw + sw/2 - Inches(0.15)
        
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, line_y - Inches(0.15), Inches(0.3), Inches(0.3))
        c.fill.solid()
        c.fill.fore_color.rgb = accent
        c.line.fill.background()
        
        step_txt = truncate(pt, 50)
        stb = slide.shapes.add_textbox(x - Inches(0.9), Inches(3.5), Inches(2), Inches(2.5))
        tf = stb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_txt
        p.font.size = Pt(11)
        p.font.color.rgb = txt
        p.alignment = PP_ALIGN.CENTER

def render_centered_statement(slide, prs, spec, theme):
    bg = hex_to_rgb(theme.color_palette.background)
    accent = hex_to_rgb(theme.color_palette.accent)
    txt = hex_to_rgb(theme.color_palette.text_on_background)
    add_bg(slide, prs, bg)
    
    title = truncate(spec.content.title, 50)
    size = fit_text(title, 35, 48, 28)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.33), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = txt
    p.alignment = PP_ALIGN.CENTER
    
    if spec.content.subtitle:
        sub = truncate(spec.content.subtitle, 70)
        sb = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.33), Inches(1))
        tf = sb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sub
        p.font.size = Pt(22)
        p.font.color.rgb = accent
        p.alignment = PP_ALIGN.CENTER

def render_comparison_columns(slide, prs, spec, theme):
    bg = hex_to_rgb(theme.color_palette.background)
    accent = hex_to_rgb(theme.color_palette.accent)
    secondary = hex_to_rgb(theme.color_palette.secondary)
    bg_txt = hex_to_rgb(theme.color_palette.text_on_background)
    sec_txt = hex_to_rgb(theme.color_palette.text_on_secondary)
    pri_txt = hex_to_rgb(theme.color_palette.text_on_primary)
    add_bg(slide, prs, bg)
    
    title = truncate(spec.content.title, 40)
    size = fit_text(title, 30, 32, 22)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = bg_txt
    
    cw = Inches(6)
    ch = Inches(5)
    add_rect(slide, Inches(0.4), Inches(1.5), cw, ch, secondary)
    add_rect(slide, Inches(6.8), Inches(1.5), cw, ch, accent)
    
    mid = max(1, len(spec.content.points) // 2)
    left_pts = spec.content.points[:mid]
    right_pts = spec.content.points[mid:]
    
    lsize = 15 if len(left_pts) > 3 else 17
    lb = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.6), Inches(4.5))
    tf = lb.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(left_pts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"-> {truncate(pt, 60)}"
        p.font.size = Pt(lsize)
        p.font.color.rgb = sec_txt
        p.space_before = Pt(10)
    
    rsize = 15 if len(right_pts) > 3 else 17
    rb = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.6), Inches(4.5))
    tf = rb.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(right_pts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"-> {truncate(pt, 60)}"
        p.font.size = Pt(rsize)
        p.font.color.rgb = pri_txt
        p.space_before = Pt(10)


RENDERERS = {
    "hero_focus": render_hero_focus,
    "split_visual_text": render_split_visual_text,
    "card_grid": render_card_grid,
    "timeline_flow": render_timeline_flow,
    "centered_statement": render_centered_statement,
    "comparison_columns": render_comparison_columns
}


def render_presentation(spec: PresentationSpec, output_dir):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    accent = hex_to_rgb(spec.theme.color_palette.accent)
    
    for s in spec.slides:
        print(f"  Slide {s.slide_id}: {s.content.title[:30]}...")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r = RENDERERS.get(s.layout.layout_family, render_split_visual_text)
        r(slide, prs, s, spec.theme)
        
        ft = slide.shapes.add_textbox(Inches(12.2), Inches(6.95), Inches(0.8), Inches(0.3))
        p = ft.text_frame.paragraphs[0]
        p.text = str(s.slide_id)
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = accent
        p.alignment = PP_ALIGN.RIGHT
    
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    fn = output_dir / f"Pro_{ts}.pptx"
    prs.save(str(fn))
    return str(fn)


class PresentationGenAgent(BaseAgent):
    name = "presentation_gen"
    description = "Create professional PowerPoint presentations with AI-generated content and images"
    icon = "📊"
    
    async def execute(self, topic: str, num_slides: int = 8, audience: str = "general", tone: str = "professional", **kwargs) -> Dict[str, Any]:
        structured_llm = self.model.with_structured_output(PresentationSpec)
        
        prompt = ChatPromptTemplate.from_messages([
            ("system", DESIGN_SYSTEM_PROMPT),
            ("human", """Create a professional presentation:
Topic: {topic}
Audience: {audience}
Tone: {tone}
Number of slides: {num_slides}

CRITICAL RULES:
1. Keep titles SHORT (max 6-8 words)
2. Keep bullet points SHORT (max 15 words each)
3. image_keyword must be ONLY ONE SINGLE WORD - a real physical object
   Good: "milkshake", "mango", "handshake", "laptop", "coffee"
   Bad: "mango milkshake glass", "happy customer", "quality product"
4. Ensure all text colors have HIGH CONTRAST with backgrounds""")
        ])
        
        chain = prompt | structured_llm
        
        print(f"\nDesigning: {topic}")
        print("="*50)
        spec = await self._safe_invoke(chain, {"topic": topic, "audience": audience, "tone": tone, "num_slides": num_slides})
        
        cp = spec.theme.color_palette
        print(f"Theme: BG={cp.background} | Accent={cp.accent}")
        print(f"Slides: {len(spec.slides)}\n")
        
        for s in spec.slides:
            kw = s.content.image_keyword or '-'
            print(f"  {s.slide_id}. [{s.layout.layout_family[:10]}] {s.content.title[:30]} [IMG: {kw}]")
        
        print("\nRendering...")
        fn = await asyncio.to_thread(render_presentation, spec, settings.PRESENTATIONS_DIR)
        print(f"\n✅ Done: {fn}")
        
        return {
            "plan": {
                "title": spec.slides[0].content.title if spec.slides else topic,
                "slides": [{"title": s.content.title, "type": s.slide_type} for s in spec.slides],
                "theme": spec.theme.model_dump()
            },
            "output_file": fn
        }
